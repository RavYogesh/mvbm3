"""Builds the four-slide executive review deck.

Kept as a script rather than a hand-made file so the deck regenerates when the
harness does. Geometry is verified by `deck/qa_deck.py`, which estimates text
extents and flags overflow and overlap -- there is no renderer in this
environment, so the layout is checked arithmetically instead of by eye.

    pip install python-pptx
    python deck/build_exec_deck.py
    python deck/qa_deck.py deck/Validation-Harness-Exec-Review.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Validation-Harness-Exec-Review.pptx"
# The `--bare` render of the published architecture page: same artefact, minus the
# title plate the slide heading already carries.
DIAGRAM = Path(__file__).resolve().parents[1] / "docs" / "architecture-diagram-bare.jpg"

# ---------------------------------------------------------------- palette
RED = RGBColor(0xD7, 0x1E, 0x28)        # Wells Fargo red
GOLD = RGBColor(0xFF, 0xCD, 0x41)       # Wells Fargo gold
INK = RGBColor(0x1C, 0x1C, 0x1C)
SLATE = RGBColor(0x5A, 0x5A, 0x5A)
MIST = RGBColor(0xEE, 0xEE, 0xEE)
LIGHT = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0x72, 0x3F)
AMBER = RGBColor(0xB3, 0x6A, 0x00)
DARK = RGBColor(0x2B, 0x0A, 0x0C)

BODY = "Arial"
W, H = 13.333, 7.5
ML = 0.72
CW = W - 2 * ML


# ---------------------------------------------------------------- helpers
def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    return box, frame


def para(frame, text, *, size=14, bold=False, color=INK, space_after=6,
         align=PP_ALIGN.LEFT, first=False, italic=False, line=None, font=BODY):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return p


def rect(slide, x, y, w, h, fill, *, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.0)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    frame = s.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.14)
    frame.margin_top = frame.margin_bottom = Inches(0.08)
    return s, frame


def chip(slide, x, y, w, text, colour, *, h=0.30, size=10):
    s, frame = rect(slide, x, y, w, h, colour, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(frame, text, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         space_after=0, first=True)
    return s


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, colour):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def header(slide, kicker, title, *, sub=None, dark=False):
    """Titles stay under ~48 characters so they hold one line at 30pt across the
    11.9in content width; qa_deck.py enforces it."""
    ink = WHITE if dark else INK
    muted = MIST if dark else SLATE
    _, f = textbox(slide, ML, 0.42, CW, 0.26)
    para(f, kicker.upper(), size=10.5, bold=True, color=GOLD if dark else RED,
         space_after=0, first=True)
    _, f2 = textbox(slide, ML, 0.74, CW, 0.62)
    para(f2, title, size=30, bold=True, color=ink, space_after=0, first=True)
    y = 1.50
    if sub:
        _, f3 = textbox(slide, ML, 1.42, CW, 0.56)
        para(f3, sub, size=13, color=muted, space_after=0, first=True, line=1.15)
        y = 2.06
    return y


def footer(slide, text, *, dark=False):
    _, f = textbox(slide, ML, 7.00, CW, 0.26)
    para(f, text, size=8.5, color=MIST if dark else SLATE, space_after=0, first=True)


# ---------------------------------------------------------------- slides
def s1_title(prs):
    s = blank(prs)
    bg(s, DARK)
    rect(s, 0, 0, W, 0.055, RED)

    _, f = textbox(s, ML, 1.78, 10.4, 0.30)
    para(f, "WELLS FARGO  ·  GENAI PLATFORM ENGINEERING", size=11, bold=True,
         color=GOLD, space_after=0, first=True)

    _, f = textbox(s, ML, 2.30, 11.4, 1.30)
    para(f, "Compressed-Model Validation Harness", size=40, bold=True, color=WHITE,
         space_after=0, first=True, line=1.05)

    _, f = textbox(s, ML, 3.86, 9.9, 0.80)
    para(f, "Turning a vendor's efficiency and quality claims into evidence a model "
            "risk committee will accept — before any model reaches our ecosystem.",
         size=15, color=MIST, space_after=0, first=True, line=1.3)

    stats = [
        ("3", "measurement suites"),
        ("12", "task families, twinned safety"),
        ("10 / 10", "instrument calibration checks"),
        ("0", "vendor-run numbers accepted"),
    ]
    for i, (value, label) in enumerate(stats):
        x = ML + i * 3.02
        _, f = textbox(s, x, 5.30, 2.85, 0.90)
        para(f, value, size=30, bold=True, color=GOLD, space_after=2, first=True)
        para(f, label, size=10.5, color=MIST, space_after=0, line=1.15)

    _, f = textbox(s, ML, 6.84, CW, 0.30)
    para(f, "Executive review  ·  August 2026  ·  Internal — pre-decisional",
         size=9.5, color=SLATE, space_after=0, first=True)


def s2_how(prs):
    """The architecture diagram itself, not a redrawn approximation of it.

    Embedding the rendered artefact keeps one source of truth: the diagram is
    regenerated from the published page by docs/export_diagram.py, so a change
    to the harness cannot leave the deck showing a stale picture. The `--bare`
    render drops the diagram's own title, which the slide heading already
    carries.
    """
    s = blank(prs)
    y = header(s, "How it works", "How a claim becomes evidence")

    if not DIAGRAM.exists():
        raise SystemExit(
            f"missing {DIAGRAM}\nrun: python docs/export_diagram.py --bare "
            f"--out docs/{DIAGRAM.name}"
        )

    # Sized from the image's real aspect ratio rather than a guessed box, so the
    # diagram is never stretched and never letterboxed.
    from PIL import Image

    with Image.open(DIAGRAM) as img:
        aspect = img.width / img.height
    height = 5.28
    width = height * aspect
    s.shapes.add_picture(str(DIAGRAM), Inches(ML), Inches(1.54),
                         width=Inches(width), height=Inches(height))

    cx = ML + width + 0.42
    cw = (ML + CW) - cx

    _, f = textbox(s, cx, 1.54, cw, 0.28)
    para(f, "READING THE FLOW", size=10.5, bold=True, color=RED, space_after=0, first=True)

    notes = [
        ("Calibration gates everything",
         "A known fault is injected first. If the harness cannot catch it, "
         "nothing below it counts as evidence."),
        ("One adapter, both models",
         "Identical prompt, decode settings and seed. The pairing is what makes "
         "the comparison valid at all."),
        ("Three suites, one decision",
         "Correctness, agentic trajectories and throughput converge on a single "
         "gate rather than three separate reports."),
    ]
    ny = 1.92
    for title, body in notes:
        _, f = textbox(s, cx, ny, cw, 1.02)
        para(f, title, size=12.5, bold=True, color=INK, space_after=4, first=True, line=1.16)
        para(f, body, size=10.5, color=SLATE, space_after=0, line=1.22)
        ny += 1.14

    _, f = textbox(s, cx, ny + 0.06, cw, 0.28)
    para(f, "AND THREE POSSIBLE ANSWERS", size=10.5, bold=True, color=RED,
         space_after=0, first=True)
    verdicts = [
        ("PASS", GREEN, "margin cleared, adequately powered"),
        ("BLOCK", RED, "past the margin, or a floor broken"),
        ("INCONCLUSIVE", AMBER, "not a soft pass"),
    ]
    vy = ny + 0.40
    for name, colour, body in verdicts:
        chip(s, cx, vy, 1.42, name, colour, h=0.27, size=8.5)
        _, f = textbox(s, cx + 1.56, vy + 0.04, cw - 1.56, 0.24)
        para(f, body, size=10, color=SLATE, space_after=0, first=True)
        vy += 0.38

    footer(s, "Interactive version: docs/architecture-diagram.html  ·  github.com/RavYogesh/mvbm3")


def s3_different(prs):
    s = blank(prs)
    y = header(s, "Why it is different", "Four decisions that carry the weight",
               sub="Each one closes a specific way a compression pitch survives an evaluation it "
                   "should have failed.")

    cards = [
        ("Non-inferiority, not averages",
         "The vendor carries the burden of proof: we assume the model IS worse until a "
         "confidence bound says otherwise. Margins are fixed before the run, so nobody "
         "can pick the threshold after seeing the result."),
        ("Cost per completed task, never per token",
         "Reasoning tokens are billed output and run near 70% of it. A model 40% cheaper "
         "per token that thinks longer and retries more is more expensive in production "
         "while winning every vendor benchmark."),
        ("Controls in code, not in the prompt",
         "A prompt is a request; a gate is an auditable control. Blocked attempts are "
         "recorded, so the report states that the model attempted an unauthorised money "
         "movement N times and the control held N times."),
        ("An answer for “we could not tell”",
         "A small study that finds no difference reads as a pass. That is the mechanism "
         "by which a degraded model gets onboarded, so the harness refuses to claim a "
         "pass it cannot defend and prints the sample size it would need."),
    ]
    cw = (CW - 0.40) / 2
    for i, (title, body) in enumerate(cards):
        x = ML + (i % 2) * (cw + 0.40)
        yy = y + (i // 2) * 1.78
        rect(s, x, yy, cw, 1.60, LIGHT)
        circle, cf = rect(s, x + 0.22, yy + 0.28, 0.52, 0.52, RED, shape=MSO_SHAPE.OVAL)
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(cf, str(i + 1), size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             space_after=0, first=True)
        _, f = textbox(s, x + 0.94, yy + 0.20, cw - 1.20, 1.30)
        para(f, title, size=14, bold=True, color=INK, space_after=5, first=True)
        para(f, body, size=11, color=SLATE, space_after=0, line=1.24)

    _, f = textbox(s, ML, y + 3.78, CW, 0.62)
    para(f, "Together these change the question from “did it score well?” to "
            "“would we have found out if it had not?” — which is the only version a "
            "model risk reviewer can sign.",
         size=12.5, bold=True, color=RED, space_after=0, first=True, line=1.25)


def s4_status(prs):
    s = blank(prs)
    bg(s, DARK)
    y = header(s, "Status and ask", "Built, tested, and calibrated", dark=True,
               sub="Running today against a simulated backend. Producing evidence needs access and "
                   "capacity, not more engineering.")

    _, f = textbox(s, ML, y, 6.30, 0.28)
    para(f, "WHAT IT ALREADY CAUGHT — IN ITS OWN FIRST VERSION", size=10.5, bold=True,
         color=GOLD, space_after=0, first=True)
    caught = [
        "A grader read the first number in a response, so any model that showed its "
        "work scored zero on correct answers.",
        "Throughput used the wrong denominator, measuring one user's experience rather "
        "than what the fleet delivers.",
        "Reasoning tokens went uncounted, understating the cost of whichever model "
        "thinks longer.",
        "Refusal was scored one-sided, so a model that refuses everything passed the "
        "strictest safety gate.",
    ]
    _, f = textbox(s, ML, y + 0.36, 6.30, 2.60)
    for i, item in enumerate(caught):
        para(f, "—  " + item, size=11.5, color=MIST, space_after=11,
             first=(i == 0), line=1.24)

    _, f = textbox(s, 7.55, y, 4.34, 0.28)
    para(f, "THE ONE CONSTRAINT", size=10.5, bold=True, color=GOLD,
         space_after=0, first=True)
    _, f = textbox(s, 7.55, y + 0.36, 4.34, 1.60)
    para(f, "66 cases cannot validate a 1-point claim", size=14, bold=True,
         color=WHITE, space_after=6, first=True, line=1.18)
    para(f, "That margin needs roughly 9,300 paired cases. It is arithmetic, not effort. "
            "Tightening the regulated margin raises the case count about ninefold — a "
            "risk-appetite call, not an engineering one.",
         size=11.5, color=MIST, space_after=0, line=1.26)

    ay = y + 2.86
    rect(s, ML, ay, CW, 0.045, RED)
    _, f = textbox(s, ML, ay + 0.28, CW, 0.30)
    para(f, "THE ASK", size=10.5, bold=True, color=GOLD, space_after=0, first=True)

    asks = [
        ("Two self-hosted endpoints", "Each candidate precision as a separate artefact, "
                                      "plus both uncompressed parents as controls."),
        ("GPU capacity for full-scale runs", "Screening is cheap. The regulated families "
                                             "are where the real capacity goes."),
        ("A decision on the regulated margin", "How much quality loss is acceptable on a "
                                               "path that feeds a customer decision."),
    ]
    aw = (CW - 0.60) / 3
    for i, (title, body) in enumerate(asks):
        x = ML + i * (aw + 0.30)
        _, f = textbox(s, x, ay + 0.70, aw, 1.10)
        para(f, title, size=13, bold=True, color=GOLD, space_after=5, first=True, line=1.16)
        para(f, body, size=11, color=MIST, space_after=0, line=1.24)

    footer(s, "29 unit tests and 10 calibration checks passing  ·  github.com/RavYogesh/mvbm3",
           dark=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    for build in (s1_title, s2_how, s3_different, s4_status):
        build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
