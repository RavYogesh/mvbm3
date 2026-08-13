"""Geometry QA for the deck, without a renderer.

This environment has no LibreOffice, so the usual render-and-look loop is not
available. Instead we estimate text extents analytically and flag the two
defects that are always user-visible: text overflowing its box, and content
crossing the slide margins.

The estimator is deliberately PESSIMISTIC (it assumes wide glyphs and generous
leading), so it over-reports rather than under-reports. A clean run is real
evidence; a flagged box needs a human look, not an automatic edit.

    python deck/qa_deck.py deck/Multiverse-Model-Validation-CTO-Review.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# Average advance width as a fraction of point size, Arial-ish. Real Arial
# lowercase averages ~0.50; we use 0.52 (0.56 bold) so mixed-case text with
# capitals and wide glyphs is not underestimated.
AVG_W = 0.52
AVG_W_BOLD = 0.56
LINE_FACTOR = 1.22          # leading multiplier
SLIDE_MARGIN_IN = 0.45      # minimum acceptable distance from slide edge

EMU_IN = 914400.0


def inches(v) -> float:
    return (v or 0) / EMU_IN


def est_lines(text: str, box_w_pt: float, size_pt: float, bold: bool) -> int:
    if not text.strip():
        return 1
    cw = size_pt * (AVG_W_BOLD if bold else AVG_W)
    if cw <= 0 or box_w_pt <= 0:
        return 1
    per_line = max(int(box_w_pt / cw), 1)
    lines = 0
    for hard in text.split("\n"):
        lines += max(1, -(-len(hard) // per_line))   # ceil div
    return lines


def frame_height_pt(shape) -> tuple[float, str]:
    """Estimated rendered height of a shape's text, in points."""
    tf = shape.text_frame
    inner_w_in = (inches(shape.width)
                  - inches(tf.margin_left) - inches(tf.margin_right))
    box_w_pt = max(inner_w_in * 72.0, 1.0)
    total = 0.0
    longest = ""
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs)
        if not p.runs:
            continue
        size = max((r.font.size.pt for r in p.runs if r.font.size), default=18.0)
        bold = any(r.font.bold for r in p.runs)
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        n = est_lines(text, box_w_pt, size, bold)
        total += n * size * LINE_FACTOR * ls
        total += (p.space_after.pt if p.space_after else 0)
        if len(text) > len(longest):
            longest = text
    total += inches(tf.margin_top) * 72 + inches(tf.margin_bottom) * 72
    return total, longest


TEXT_BOX = 17   # MSO_SHAPE_TYPE.TEXT_BOX


def overlaps(slide, idx: int) -> list[str]:
    """Plain text boxes must never intersect each other.

    Only text boxes are compared. Filled rectangles are deliberate backdrops
    that text sits on top of, so including them would flag every card on every
    slide and drown the real collisions.
    """
    boxes = []
    for sh in slide.shapes:
        if sh.shape_type != TEXT_BOX or not sh.has_text_frame:
            continue
        if not sh.text_frame.text.strip():
            continue
        boxes.append((inches(sh.left), inches(sh.top), inches(sh.width),
                      inches(sh.height), sh.text_frame.text[:44].replace("\n", " / ")))
    out = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
            oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
            if ox > 0.05 and oy > 0.05:
                out.append(f"  slide {idx:>2}: OVERLAP {oy:.2f}in vertical — "
                           f"\"{a[4]}\" x \"{b[4]}\"")
    return out


def main(path: str) -> int:
    prs = Presentation(path)
    sw, sh = inches(prs.slide_width), inches(prs.slide_height)
    problems: list[str] = []
    checked = 0

    for idx, slide in enumerate(prs.slides, 1):
        problems.extend(overlaps(slide, idx))
        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)

            # Off-slide / margin encroachment. Full-bleed bars are intentional
            # and excluded by checking only shapes that are not slide-wide.
            if w < sw - 0.01:
                if x < SLIDE_MARGIN_IN - 0.001 or x + w > sw - SLIDE_MARGIN_IN + 0.001:
                    problems.append(
                        f"  slide {idx:>2}: horizontal margin — x={x:.2f} w={w:.2f} "
                        f"(right edge {x+w:.2f}, slide {sw:.2f})")
            if y < 0 or y + h > sh + 0.01:
                problems.append(f"  slide {idx:>2}: vertical overflow — y={y:.2f} h={h:.2f}")

            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            checked += 1
            need_pt, longest = frame_height_pt(shape)
            have_pt = h * 72.0
            # 6pt of slack: our estimator rounds up, and PowerPoint's autofit is
            # off, so a hair over is not a real defect.
            if need_pt > have_pt + 6.0:
                problems.append(
                    f"  slide {idx:>2}: TEXT OVERFLOW — needs ~{need_pt/72:.2f}in, "
                    f"box is {h:.2f}in  @({x:.2f},{y:.2f}) w={w:.2f}\n"
                    f"           \"{longest[:78]}\"")

    print(f"\nGeometry QA — {path}")
    print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} slides, {checked} text frames checked\n")
    if problems:
        print("\n".join(problems))
        print(f"\n  {len(problems)} potential issue(s).\n")
        return 1
    print("  No overflow or margin issues detected.\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1] if len(sys.argv) > 1
                          else str(Path(__file__).parent /
                                   "Multiverse-Model-Validation-CTO-Review.pptx")))
